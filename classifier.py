"""
Text Classifier
===============
Uses Hugging Face zero-shot classification to label text documents
without any fine-tuning or labeled training data.
"""

import logging
from functools import lru_cache
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

# Lazy imports — only pulled in when classify() is first called
_pipeline = None


def _get_pipeline(model_name: str):
    """Load (and cache) the HuggingFace zero-shot classification pipeline."""
    global _pipeline
    if _pipeline is None:
        try:
            from transformers import pipeline
            logger.info(f"Loading NLP model '{model_name}' (first run may download weights)…")
            _pipeline = pipeline("zero-shot-classification", model=model_name)
            logger.info("NLP model loaded.")
        except ImportError:
            logger.error("transformers library not installed. Run: pip install transformers")
            raise
    return _pipeline


def extract_text(file_path: Path) -> Optional[str]:
    """
    Extract plain text from a file.
    Supports: .txt, .md, .html, .json, .xml, .csv, .pdf, .docx, .pptx, .xlsx
    Returns None if extraction fails or file type is unsupported.
    """
    suffix = file_path.suffix.lower()

    try:
        if suffix in {".txt", ".md", ".csv", ".json", ".xml", ".html"}:
            return file_path.read_text(errors="ignore")[:4000]

        if suffix == ".pdf":
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                text = "\n".join(p.extract_text() or "" for p in pdf.pages[:5])
            return text[:4000] or None

        if suffix in {".docx", ".doc"}:
            import docx
            doc = docx.Document(file_path)
            text = "\n".join(p.text for p in doc.paragraphs)
            return text[:4000] or None

        if suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return "\n".join(texts)[:4000] or None

        if suffix == ".xlsx":
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            texts = []
            for ws in wb.worksheets[:3]:
                for row in ws.iter_rows(max_row=20, values_only=True):
                    texts.append(" ".join(str(c) for c in row if c is not None))
            return "\n".join(texts)[:4000] or None

    except Exception as e:
        logger.warning(f"Text extraction failed for {file_path.name}: {e}")

    return None


def classify(
    text: str,
    candidate_labels: list[str],
    model_name: str,
    confidence_threshold: float = 0.4,
) -> tuple[str, float]:
    """
    Run zero-shot classification on text.

    Returns:
        (label, score) — best label and its confidence score.
        Falls back to ("unknown", 0.0) if below threshold.
    """
    if not text or not text.strip():
        return "unknown", 0.0

    clf = _get_pipeline(model_name)
    result = clf(text[:1024], candidate_labels=candidate_labels, multi_label=False)

    best_label: str = result["labels"][0]
    best_score: float = result["scores"][0]

    if best_score < confidence_threshold:
        logger.debug(f"Low confidence ({best_score:.2f}) → labelled 'unknown'")
        return "unknown", best_score

    logger.debug(f"Classified as '{best_label}' ({best_score:.2f})")
    return best_label, best_score
